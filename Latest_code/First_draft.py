# -*- coding: utf-8 -*-
"""
Created on Mon Jan 20 12:36:08 2020
@author: sbiswas149
"""

#import win32com.client
import json
import requests
import shutil
import mammoth
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import csv
import xlrd
import re
import pandas as pd
import openpyxl
from openpyxl.styles import Font, PatternFill
from pyxlsb import open_workbook
from independentsoft.msg import Message
from striprtf.striprtf import rtf_to_text
from odf import text
from odf.opendocument import load
from sqlalchemy import create_engine
import configparser
import textract
from pdf2image import convert_from_path, convert_from_bytes
import logging.config
import os, sys
from PIL import Image
#from visio2pdf import Visio2PDFConverter
import datetime
import botocore
from boto3.session import Session
import boto3
from botocore.exceptions import NoCredentialsError, ClientError

outputCol = 0
outputRow = 1
memory_error_files = []
exception_files = []
keyword_dic = {}
pattern_list = {}
raw_details_list = []
try:
    config = configparser.ConfigParser()
    config.read('properties.ini')
    config_folder = config['Directories']['config_file']
    output_file = config['Directories']['output_file']
    temp_storage_path = config['Directories']['temp_storage']
    address = config['Mysql']['address']
    user = config['Mysql']['user']
    password = config['Mysql']['password']
    database_name = config['Mysql']['database_name']
    table_name = config['Mysql']['table_name']
    secret_key = config['AWS']['secret_key']
    access_key = config['AWS']['access_key']
except KeyError as e:
    print("Config file is empty")
    print('Config file is empty - fatal error')
    sys.exit(1)
except Exception as e:
    print("Check config file", e)
    sys.exit(1)

# create logger
logger = logging.getLogger('VDR')
logging.basicConfig(filename= config_folder,
                                filemode='a',
                                format='%(asctime)s,%(msecs)d %(name)s %(levelname)s %(message)s',
                                datefmt='%H:%M:%S',
                                level=logging.DEBUG)
logging.warning('is when this event was logged.')

# converting numeric value to character equivalent to excel header
# value 2 equivalent to char B
def cell_header(a):
    if a <= 26:
        return chr(a + 64)
    count = a // 26 + 64
    rest = a % 26 + 64
    return chr(count) + chr(rest)


# processes csv file
def csv_handler(path):
    with open(path) as csvfile:
        readCSV = csv.reader(csvfile)
        # count is row number
        count = 0
        page_count = 1
        for row in readCSV:
            count = count + 1
            for i in range(len(row)):
                # print(row[i] + " " + str(count) + " " + str(i))
                # row is representing each row  of the csv file - a list
                list_details = folderName(path)
                cell_details = cell_header(i + 1) + str(count)
                sheet_handler(list_details, cell_details, path, page_count, row[i])

# checks whether a value is nan
def isNan(val):
    return val != val

# processes excel file
def xlsx_handler(path):
    try:
        sheets_dict = pd.read_excel(path, sheet_name=None, header=None)
        sheet_count = len(sheets_dict)
        list_details = folderName(path)
        for name, sheet in sheets_dict.items():
            # removing columns which are in float format
            #sheet = sheet.loc[:, sheet.dtypes != np.float64]
            for index, row in sheet.iterrows():
                for i in range(len(row)):
                    cell = cell_header(i + 1) + str(index + 1)
                    cell_details = name + ": " + cell
                    sheet_handler(list_details, cell_details, path, sheet_count, row[i])
    except Exception as e:
        print("Issue processing xlsx::", e)
        pass

# method which writes all the output data to a list of dictionary - raw_details_list
def sheet_handler(list_details, cell_details, path, page_count, row):
    try:
        global keyword_dic, pattern_list
        if isinstance(row, str):
            for category, list_of_pattern in pattern_list.items():
                for index in range(len(list_of_pattern)):
                    words = list_of_pattern[index].findall(row)
                    for i in words:
                        temp_dic = {"File Name": list_details[0], "Index Number" : list_details[1],
                                    "Folder Name": list_details[2], "File Path" : path, "File Type": list_details[3],
                                    "Total Number of Page(s)/Slide(s)": page_count, "Keyword": keyword_dic[category][index], "IT Category":
                                    category, "Page/Slide Number": cell_details}
                        raw_details_list.append(temp_dic)

    except Exception as e:
        print("Error while writing to RAW details", e)
        sys.exit(1)

# removes unneccessary spaces from extracted texts from images
def removal(text):
    return text.replace("\\r", " ").replace("\\n", " ").replace('\\x', " ").replace('0c', " ")

# handles these extensions - ".PNG", ".png", ".JPEG", ".jpeg", ".JPG", ".jpg", ".gif", ".pnm", ".PNM"
def img_handler(path):
    list_details = folderName(path)
    text = textract.process(path, encoding='ascii', method='tesseract')
    content = removal(str(text))
    sheet_handler(list_details, 1, path, 1, content)

# handles these extensions - ".tiff", ".tif", ".jfif", ".bmp"
def tiff_handler(path):
    list_details = folderName(path)
    out_file = os.path.join(temp_storage_path, "temp_img.JPEG")
    im = Image.open(path)
    out = im.convert("RGB")
    out.save(out_file, "JPEG", quality=90)
    img_handler_for_pdf(out_file, path, 1, 1, list_details)
    os.remove(out_file)

def write_image(shape, count, list_details, total_pages, current_page):
    image = shape.image
    image_bytes = image.blob
    image_filename = 'image{:03d}.{}'.format(count, image.ext)
    with open(image_filename, 'wb') as fi:
        fi.write(image_bytes)
        fi.close()
    #txt = pytesseract.image_to_string(PILImage.open(image_filename))
    txt = textract.process(image_filename, encoding='cp437', method='tesseract')
    sheet_handler(list_details, current_page, path, total_pages, txt)
    os.remove(image_filename)

# converts ppt to pdf for further processing
def ppt_handler(path):
    prs = Presentation(path)
    total_pages = len(prs.slides)
    list_details = folderName(path)
    page = 0
    for slide in prs.slides:
        page = page + 1
        count = 0
        for shape in slide.shapes:
            count = count + 1
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                write_image(shape, page, list_details, total_pages, page)
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        #print(run.text)
                        sheet_handler(list_details, page, path, total_pages, run.text)

# coverts the docx to pdf for further processing
def doc_handler(path):
    text = ""
    list_details = folderName(path)
    with open(path, "rb") as docx_file:
        result = mammoth.extract_raw_text(docx_file)
        text += result.value
    sheet_handler(list_details, 1, path, 1, text)

# extracts images out of scanned pdfs
def pdf_handler(path):
    try:
        pages = convert_from_path(path, 80)
        total_pages = len(pages)
        list_details = folderName(path)
        for index, page in enumerate(pages):
            temp_file = list_details[0] + str(index + 1) + '.jpeg'
            temp_path = os.path.join(temp_storage_path, temp_file)
            page.save(temp_path, 'JPEG')
            img_handler_for_pdf(temp_path, path, index, total_pages, list_details)
            os.remove(temp_path)

    except Exception as e:
        print("exception has occured while procesing this file:: ", path)
        print("The error is:: ", e)
        print("skipping this file")
        pass

# handles images converted from pdfs which in turn is converted from docx and ppts
def img_handler_for_pdf(temp_path, path, page_number ,total_pages, list_details):
    try:
        text = textract.process(temp_path, encoding="ascii",  errors='ignore', method='tesseract')
        #text = tools[0].image_to_string(Image.open(temp_path), lang="eng", builder=pyocr.builders.TextBuilder())
        #text = pytesseract.image_to_string(temp_path)
        content = removal(str(text))
        sheet_handler(list_details, page_number, path, total_pages, content)
    except Exception as e:
        print("exception has occured while procesing this file:: ", path)
        print("The error is:: ", e)
        print("skipping this file")
        pass

# handles xlsb files
def xlsb_handler(path):
    sheet_count = len(open_workbook(path).sheets)
    list_details = folderName(path)
    with open_workbook(path) as wb:
        for sheetname in wb.sheets:
            row_number = 0
            with wb.get_sheet(sheetname) as sheet:
                for row in sheet.rows():
                    row_number = row_number + 1
                    for i in range(len(row)):
                        if row[i].v:
                            cell = cell_header(i + 1) + str(row_number)
                            cell_details = sheetname + ": " + cell
                            sheet_handler(list_details, cell_details, path, sheet_count, row[i].v)


# odt file handler
def odt_handler(path):
    textdoc = load(path)
    allparas = textdoc.getElementsByType(text.P)
    list_details = folderName(path)
    for i in allparas:
        sheet_handler(list_details, 1, path, 1, str(i))

# processes msg files
def msg_handler(path):
    content = str(Message(path).body)
    sub = str(Message(path).subject)
    list_details = folderName(path)
    sheet_handler(list_details, 1, path, 1, content)
    sheet_handler(list_details, 1, path, 1, sub)

# processes text files
def txt_handler(path):
    list_details = folderName(path)
    file = open(path, mode='r')
    content = file.read()
    sheet_handler(list_details, 1, path, 1, content)
    file.close()

#processes rtf files
def rtf_handler(path):
    list_details = folderName(path)
    file = open(path, mode='r')
    rtf = file.read()
    content = rtf_to_text(rtf)
    sheet_handler(list_details, 1, path, 1, content)
    file.close()

# extracts index number from file name
def index_number_verification(index):
    r = re.compile('[@_!#$%^&*()<>?/\|}{~:]')
    if "_" in index:
        index = index.split("_")[0]
    if (index.isupper() or index.islower()) or r.search(index) != None:
        return " "
    index = re.sub('[^0-9\.\s]+', '', index)
    if len(index) == 1:
        index = re.sub('[.]', " ", index)
    return index

# generates file details
def folderName(path):
    # C:\Users\sb512911\Desktop\All\Screenshots\1.2.1 ILS Structure Chart.csv
    path = os.path.normpath(path)
    file_name = path.split(os.sep)[-1]
    index = ""
    list1 = ['1','2','3','4','5','6','7','8','9','0', '.', '_']
    for i in file_name:
        if i in list1:
            index = index + i
        else:
            break
    index_number = index_number_verification(index)
    folder_name = path.split(file_name)[0]
    file_type = path.split(".")[-1]
    return [file_name, index_number, folder_name, file_type]


# pass the additional parameter for relevance test - dynamically
# generates file level details sheet from "Raw Details"
def generatingFileLevelDetail(file_path):
    try:
        data = pd.read_excel(open(file_path, 'rb'), sheet_name=0)
        if not data.empty:
            temp_data = data[["File Name", "Index Number", "Folder Name", "File Type"]]
            temp_data["Scan Status"] = "Successful"
            file_level_detail = temp_data[temp_data['File Type'] != "Unknown"] \
                .drop_duplicates().sort_values('File Name')

            file_level_detail["Number of Keywords found"] = \
                data[data["File Type"] != "Unknown"][['File Name']].sort_values('File Name').groupby(
                    ['File Name']).size().reset_index()[0].values

            file_level_detail["Number of Keyword Categories"] = \
                data[data["File Type"] != "Unknown"][['File Name', 'IT Category']].drop_duplicates().groupby(
                    ['File Name']).size() \
                    .reset_index()[0].values

            intermediate1 = data[data["File Type"] != "Unknown"][['File Name', 'Keyword']].groupby(
                ['File Name', 'Keyword']).size().reset_index()
            intermediate1[[0]] = intermediate1[[0]].astype(str).values
            intermediate1['joined'] = intermediate1[['Keyword', 0]].apply(lambda x: '('.join(x) + "), ", axis=1).values
            file_level_detail["Keyword Details"] = intermediate1[["File Name", "joined"]].groupby(
                ["File Name"]).sum().values

            file_level_detail["Category Details"] = \
                data[data["File Type"] != "Unknown"].groupby("File Name").apply(
                    lambda x: x["IT Category"].drop_duplicates().str.cat(sep=", ")).values

            file_level_detail["File Path"] = data[data["File Type"] != "Unknown"]["File Path"].drop_duplicates().values

            # Relevance check -- code here - assuming relevant count = 5
            mapping = {True: "Yes", False: "No"}
            file_level_detail["Relevant?"] = \
            data[data["File Type"] != "Unknown"][["File Name", "Keyword"]].drop_duplicates()[["File Name"]].groupby(
                ["File Name"]).size(). \
                reset_index()[0].apply(lambda x: x > 5).map(mapping).values

            temp_data_unknown = temp_data[["File Name", "Index Number", "Folder Name", "File Type"]][
                temp_data['File Type'] == "Unknown"]
            temp_dict = {"Number of Keywords found": 0, "Number of Keyword Categories": 0, "Keyword Details": "", \
                         "Category Details": "", "Relevant?": "", "File Path": "", "Scan Status": "Unsuccessful"}
            temp_df = pd.DataFrame([temp_dict])
            unknown_files = cross_join(temp_data_unknown, temp_df)
            file_level_detail = file_level_detail.append(unknown_files, ignore_index=True)
            writing_to_excel(file_level_detail, "File Level Detail", "a")

    except:
        print("error occured while generating file level detail")
        sys.exit(1)


# generates keyword detail output sheet sheet
def generatingKeywordDetails(file_path):
    data = pd.read_excel(open(file_path, 'rb'), sheet_name="Raw Details")
    try:
        if not data.empty:
            # distinct keywords and categories
            keyword_details = data[["Keyword", "IT Category"]].drop_duplicates()

            #  Number of files with keywords
            keyword_details = data[["Keyword", "File Name"]].drop_duplicates().groupby(["Keyword"]) \
                .size().reset_index().merge(keyword_details, on="Keyword")

            # renaming
            keyword_details = keyword_details. \
                rename(columns={'IT Category': 'Keyword Category', 0: 'Number of files with keywords'})

            # Total keyword hits
            keyword_details = data[["File Name", "Keyword"]].groupby(["Keyword"]) \
                .size().reset_index().merge(keyword_details, on="Keyword").rename(columns={0: "Total keyword hits"})

            # File(s) with maximum keyword hits and File Path(s)
            intermediate2 = data[["Keyword", "File Name", "File Path"]].groupby(["Keyword", "File Name", "File Path"]) \
                .size().reset_index().merge(keyword_details, on="Keyword").rename({0: "Count"}, axis=1) \
                [["File Name", "Keyword", "File Path", "Count"]]

            intermediate3 = intermediate2.groupby('Keyword')['Count'].apply(lambda x: x.eq(x.max()))

            intermediate4 = intermediate2.loc[intermediate3].groupby(['Keyword'])['File Name'].agg(
                ', '.join).reset_index()

            intermediate5 = intermediate2.loc[intermediate3].groupby(['Keyword'])['File Path'].agg(
                ', '.join).reset_index()

            keyword_details = keyword_details.merge(intermediate4, on="Keyword")

            keyword_details = keyword_details.merge(intermediate5, on="Keyword")

            keyword_details = keyword_details.rename(columns={"File Path": "File Path(s) with maximum keyword", \
                                                              "File(s) with maximum keyword hits": "File Path(s) with maximum keyword"})

            # All file(s) with keyword hits
            keyword_details["File Path(s)"] = \
                data[["Keyword", "File Path"]].groupby(["Keyword"]) \
                    .apply(lambda x: x["File Path"].drop_duplicates().str.cat(sep=", ")).reset_index() \
                    .merge(keyword_details, on="Keyword")[[0]].values

            # Index Number(s)
            keyword_details["Index Number(s)"] = \
                data[["Keyword", "Index Number"]].groupby("Keyword") \
                    .apply(lambda x: x["Index Number"].drop_duplicates().astype(str).str.cat(sep=", ")).reset_index() \
                    .merge(keyword_details, on="Keyword")[[0]].values

            # All the file names which has the keyword
            keyword_details["File Name(s)"] = \
                data[["Keyword", "File Name"]].groupby("Keyword") \
                    .apply(lambda x: x["File Name"].drop_duplicates().str.cat(sep=", ")).reset_index() \
                    .merge(keyword_details, on="Keyword")[[0]].values

            keyword_details = keyword_details.rename(columns={"File Name": "File Name(s) with maximum keyword", \
                                                              "File Name(s)": "All file(s) with keyword hits"})

            keyword_details = keyword_details[["Keyword", "Keyword Category", "Number of files with keywords", \
                                               "Total keyword hits", "File Name(s) with maximum keyword",
                                               "File Path(s) with maximum keyword", \
                                               "All file(s) with keyword hits", "Index Number(s)", "File Path(s)"]]

            # writing to "File Level Detail" sheet
            writing_to_excel(keyword_details, "Keyword Details", "a")

    except Exception as e:
        print("Error occured while generating Keyword details", e)
        sys.exit(1)

# handles DRL file and writes to DRL Details sheet
def generatingDRLDetails(drl_path, output_path):
    df = pd.read_excel(output_path, sheet_name="Raw Details")
    try:
        if not df.empty:
            data = df[df["File Type"] != "Unknown"][["File Name", "Index Number", "File Path", "Keyword"]].drop_duplicates()
            dict_file_path = {}
            all_files = data["File Name"].tolist()
            dict_all = {}
            dict_index_number = {}
            each_row = []
            for file in all_files:
                dict_all[file] = data[data['File Name'] == file]['Keyword'].tolist()
                dict_index_number[file] = data[data['File Name'] == file]['Index Number'].tolist()[0]
                dict_file_path[file] = data[data['File Name'] == file]['File Path'].tolist()[0]
            workbook1 = xlrd.open_workbook(drl_path)
            sh = workbook1.sheet_by_index(0)
            for rownum in range(sh.nrows):
                rows = sh.row_values(rownum)
                if rows[5].strip().lower() == "it":
                    sentence = rows[6]
                    drl_number = int(rows[1])
                    satisfied_output_file_name = {}
                    satisfied_output_file_path = {}
                    satisfied_output_index = {}
                    flag = "No"
                    relevant = 0
                    # keys = file name
                    for keys, values in dict_all.items():
                        count = 0
                        for m in values:
                            r1 = re.compile('|'.join([r'(?<!\w)%s(?!\w)' % re.escape(m)]), flags=re.I)
                            matches = r1.findall(sentence)
                            if matches:
                                count = count + 1
                        if count >= 5:
                            flag = "Yes"
                            relevant = relevant + 1
                            if drl_number not in satisfied_output_file_name:
                                satisfied_output_file_name[int(drl_number)] = keys + ', '
                                satisfied_output_file_path[int(drl_number)] = dict_file_path[keys] + ", "
                                satisfied_output_index[int(drl_number)] = str(dict_index_number[keys]) + ", "
                            else:
                                satisfied_output_file_name[int(drl_number)] += keys + ", "
                                satisfied_output_file_path[int(drl_number)] += dict_file_path[keys] + ", "
                                satisfied_output_index[int(drl_number)] += str(dict_index_number[keys]) + ", "
                    if (satisfied_output_file_name and satisfied_output_file_path \
                            and satisfied_output_index):
                        dict1 = {}
                        dict1["DRL #"] = drl_number
                        dict1["Request"] = sentence
                        dict1["Relevant Files Found?"] = flag
                        dict1["Number of relevant files"] = relevant
                        dict1["File Name(s)"] = satisfied_output_file_name[drl_number]
                        dict1["Index Number(s)"] = satisfied_output_index[drl_number]
                        dict1["File Path(s)"] = satisfied_output_file_path[drl_number]
                        each_row.append(dict1)
                    else:
                        dict1 = {}
                        dict1["DRL #"] = drl_number
                        dict1["Request"] = sentence
                        dict1["Relevant Files Found?"] = flag
                        dict1["Number of relevant files"] = relevant
                        dict1["Index Number(s)"] = ""
                        dict1["File Path(s)"] = ""
                        dict1["File Name(s)"] = ""
                        each_row.append(dict1)

            dataframe = pd.DataFrame(each_row)
            dataframe = dataframe[
                    ["DRL #", "Request", "Relevant Files Found?", "Number of relevant files", "File Name(s)", "Index Number(s)",
                     "File Path(s)"]]
            # writing to "File Level Detail" sheet
            writing_to_excel(dataframe, "DRL Details", "a")
    except Exception as e:
        print("Error has occured while generating DRL details:: ", e)

# writes, formats, colours the final output file
def writing_to_excel(dataframe, sheet_name, mode):
    if not os.path.exists(output_file):
        open(output_file, 'w').close()
    with pd.ExcelWriter(output_file, engine='openpyxl', mode=mode) as writer:
        dataframe.to_excel(writer, sheet_name, index=False)
        writer.save()

    wb1 = openpyxl.load_workbook(output_file)
    ws = wb1[sheet_name]
    fillBack = PatternFill(start_color="00FFFF00", fill_type="solid")

    ws.column_dimensions['A'].width = 24.94
    ws.column_dimensions['B'].width = 24.94
    ws.column_dimensions['C'].width = 24.94
    ws.column_dimensions['D'].width = 24.94
    ws.column_dimensions['E'].width = 24.94
    ws.column_dimensions['F'].width = 24.94
    ws.column_dimensions['G'].width = 24.94
    ws.row_dimensions[1].height = 28.8

    for cell in ws["1:1"]:
        cell.fill = fillBack
    wb1.save(output_file)

# does a cross join in python
def cross_join(left, right):
    return left.assign(key=1).merge(right.assign(key=1), on='key').drop('key', 1)

# creates regex pattern for each keyword
def pattern(keys):
    return re.compile('|'.join([r'(?<!\w)%s(?!\w)' % re.escape(keys)]), flags=re.I)

# creating a dictionary of lists containing the regex pattern for each keyword
def keywordlist():
    global keyword_dic, pattern_list
    df = keyword_extraction()
    category_list = df["category"].drop_duplicates().tolist()
    for i in category_list:
        temp = df[df["category"] == i]["keyword"].tolist()
        keyword_dic[i] = temp
        pattern_list[i] = [pattern(key) for key in temp]

# extracts the data dictionary from mysql table - data_dictionary.virtual_data_room
def keyword_extraction():
    try:
        #db_connection_str = 'mysql+pymysql://root:admin@localhost/virtual_data_room'
        db_connection_str = 'mysql+pymysql://' + user + ":" + password + "@" + address + "/" + database_name
        db_connection = create_engine(db_connection_str)
        query = 'SELECT keyword, category FROM {}'.format(table_name)
        df = pd.read_sql(query, con=db_connection)
        return df.drop_duplicates()
    except:
        print("mysql connection issue")
        sys.exit(1)

def multi_processing(i):
    print("file name ", i)
    if i.endswith(".csv"):
        csv_handler(i)
    elif i.endswith((".xlsb")):
        xlsb_handler(i)
    elif i.endswith(".msg"):
        msg_handler(i)
    elif i.endswith(".txt"):
        txt_handler(i)
    elif i.endswith(".rtf"):
        rtf_handler(i)
    elif i.endswith((".docx", ".doc")):
        print(datetime.datetime.now())
        doc_handler(i)
    elif i.endswith(".odt"):
        odt_handler(i)
    elif i.endswith((".pptx", "ppt")):
        print(datetime.datetime.now())
        ppt_handler(i)
    elif i.endswith((".PNG", ".png", ".JPEG", ".jpeg", ".JPG", ".jpg", ".gif", ".pnm", ".PNM")):
        img_handler(i)
    elif i.endswith((".xlsx", ".xlsm", ".xls", ".XLS", "XLSX")):
        print(datetime.datetime.now())
        xlsx_handler(i)
    elif i.endswith((".pdf", ".PDF")):
        print(datetime.datetime.now())
        pdf_handler(i)
    elif i.endswith((".tiff", ".tif", ".jfif", ".bmp")):
        tiff_handler(i)
    else:
        # files which are having extension other than the above ones, will be declared Unknown
        list_details = folderName(i)
        temp_dic = {"File Name": list_details[0], "Index Number": list_details[1],
                    "Folder Name": list_details[2], "File Path": i, "File Type": "Unknown",
                    "Total Number of Page(s)/Slide(s)": "", "Keyword": "", "IT Category":
                    "", "Page/Slide Number": ""}
        raw_details_list.append(temp_dic)

def file_details(s3_directory, guid, user_folder, folder):
    bucket_name = s3_directory.split("//")[1].split("/")[0]
    try:
        session = Session(aws_access_key_id=access_key,
                          aws_secret_access_key=secret_key)
        s3 = session.resource('s3')
        your_bucket = s3.Bucket(bucket_name)
        if not os.path.exists(user_folder):
            os.makedirs(user_folder, exist_ok=True)
        for s3_file in your_bucket.objects.all():
            if s3_file.key.startswith(guid + folder):
                path, filename = os.path.split(s3_file.key)
                path_local = user_folder + path
                if not os.path.exists(path_local):
                    os.makedirs(path_local)
                os.chdir(path_local)
                s3.meta.client.download_file(bucket_name, s3_file.key, filename)
        return user_folder + "/" + guid + folder

    except botocore.exceptions.ClientError as e:
        if e.response['Error']['Code'] == "404":
            print("The object does not exist.")
        else:
            raise

def upload_to_s3(output_file, bucket, object_name):
    session = Session(aws_access_key_id=access_key,
                      aws_secret_access_key=secret_key)
    s3 = session.resource('s3')
    try:
        s3.meta.client.upload_file(output_file, bucket, object_name)
    except ClientError as e:
        logging.error(e)
        return False
    return True


if __name__ == '__main__':
    try:
        # intermediate file is the path where we have these details - VDR path, DRL path, additional
        root = sys.argv[1]            #s3://bucketname/2222/VDR
        drl_folder = sys.argv[2]      #s3://bucketname/2222/DRL
        keywords = sys.argv[3]        #"abc1,abc2"
        guid = sys.argv[4]            #2222
        drl_file_path = ""
        vdr = file_details(root, guid, user_folder = "/home/ddvadmin/scripts/File/", folder = "/VDR")

        drl_folder = file_details(drl_folder, guid, user_folder = "/home/ddvadmin/scripts/File/", folder = "/DRL")
        for path, subdirs, files in os.walk(drl_folder):
            for name in files:
                a = os.path.join(path, name)
                drl_file_path = a

        additional_keywords = [i.strip() for i in keywords.split(",")]

        keywordlist()

        # adding the user given keywords
        keyword_dic["User Input"] = [pattern(key) for key in additional_keywords]

        fileList = []
        for path, subdirs, files in os.walk(vdr):
            for name in files:
                a = os.path.join(path, name)
                fileList.append(a)

        for file in fileList:
            multi_processing(file)
        shutil.rmtree(vdr)
        raw_dataframe = pd.DataFrame(raw_details_list)
        writing_to_excel(raw_dataframe, "Raw Details", "o")

        generatingFileLevelDetail(output_file)
        generatingKeywordDetails(output_file)
        generatingDRLDetails(drl_file_path, output_file)

        out = guid + "/output/output.xlsx"
        bucket_name = root.split("//")[1].split("/")[0]
        print(output_file)
        if upload_to_s3(output_file, bucket_name, out):
            data_set = {"url": out, "guid": guid}
            json_dump = json.dumps(data_set)
            url = 'http://172.30.26.82/ProcessEmailRequest/'
            headers = {'content-type': 'application/json', 'Accept-Charset': 'UTF-8'}
            result = requests.post(url, data=json_dump, headers=headers)

    except KeyError as e:
        print("Either data is not present or column is missing", e)
        logger.critical("Either data is not present or column is missing", e)
    except PermissionError as e:
        print("Some files are open or not have permission to open - thus not accessible", e)
    except Exception as e:
        logger.info("Fatal error has occured - ", e)
    finally:
        print(datetime.datetime.now())
        logger.info("The end")
        logger.info("***********************************NEXT RUN************************************")
        logger.info("*******************************************************************************")
        logger.info("*******************************************************************************")
