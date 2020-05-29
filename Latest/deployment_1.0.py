# -*- coding: utf-8 -*-
"""
Created on Mon Jan 20 12:36:08 2020
@author: sbiswas149
"""
from io import StringIO
from bs4 import BeautifulSoup
from tika import parser
import time
import os
import sys
import json
import requests
import shutil
import subprocess
import mammoth
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import csv
import xlrd
import re
import pandas as pd
import openpyxl
from openpyxl.styles import Font, PatternFill
from pyxlsb import open_workbook
from independentsoft.msg import Message
from striprtf.striprtf import rtf_to_text
from odf import text
from odf.opendocument import load
from sqlalchemy import create_engine
import configparser
import textract
from pdf2image import convert_from_path, convert_from_bytes
import logging.config
from PIL import Image
import datetime
import botocore
from boto3.session import Session
import numpy as np
from botocore.exceptions import NoCredentialsError, ClientError
from http import HTTPStatus
import pytesseract


keyword_dic = {}
pattern_list = {}
raw_details_list = []
warning_message = ""
file_skip_count = 0
skipped_file_name = []
unknown_file_extension = []
files_with_matching_keywords = []
all_files = []
try:
    config = configparser.ConfigParser()
    config.read('properties.ini')
    log_dir = config['Directories']['log_file']
    # output_file = config['Directories']['output_file']
    temp_storage_path = config['Directories']['temp_storage']
    user_folder = config['Directories']['user_folder']
    address = config['Mysql']['address']
    user = config['Mysql']['user']
    password = config['Mysql']['password']
    database_name = config['Mysql']['database_name']
    table_name = config['Mysql']['table_name']
    secret_key = config['AWS']['secret_key']
    access_key = config['AWS']['access_key']
except KeyError as e:
    #logger.info("config file is empty")
    sys.stdout.write("Config file is empty \n")
    sys.stdout.write('Config file is empty - fatal error, check the error: {} \n'.format(e))
    sys.exit(1)
except Exception as e:
    sys.stdout.write("Check config file, this error has occurred: {} \n".format(e))
    sys.exit(1)

# create logger
timestr = str(time.strftime("%Y%m%d-%H%M%S")) + ".logs"
log_file = os.path.join(log_dir, timestr)
logger = logging.getLogger('VDR')
logging.basicConfig(filename=log_file, filemode='a',
                    format='%(asctime)s,%(name)s %(levelname)s %(message)s',
                    datefmt='%H:%M:%S', level=logging.DEBUG)
logging.warning('is when this event was logged.')


# converting numeric value to character equivalent to excel header
# value 2 equivalent to char B
def cell_header(a):
    if a <= 26:
        return chr(a + 64)
    count = a // 26 + 64
    rest = a % 26 + 64
    return chr(count) + chr(rest)


# processes csv file
def csv_handler(path):
    global file_skip_count
    list_details = folderName(path)
    try:
        with open(path) as csvfile:
            readCSV = csv.reader(csvfile)
            # count is row number
            count = 0
            page_count = 1
            for row in readCSV:
                count = count + 1
                for i in range(len(row)):
                    # print(row[i] + " " + str(count) + " " + str(i))
                    # row is representing each row  of the csv file - a list
                    cell_details = cell_header(i + 1) + str(count)
                    sheet_handler(list_details, cell_details, path, page_count, row[i])
    except Exception as e:
        file_skip_count += 1
        skipped_file_name.append({"File Name": list_details[0]})
        logging.error("Error has occured while processing this file {} and the error is {}".format(list_details[0],e))
        sys.stdout.write('Error has occurred while processing this file: {}, error is {} \n'.format(list_details[0], e))
        pass


# processes excel file
def xlsx_handler(path):
    global file_skip_count
    list_details = folderName(path)
    try:
        sheets_dict = pd.read_excel(path, sheet_name=None, header=None)
        sheet_count = len(sheets_dict)
        for name, sheet in sheets_dict.items():
            # removing columns which are in float format
            # sheet = sheet.loc[:, sheet.dtypes != np.float64]
            for index, row in sheet.iterrows():
                for i in range(len(row)):
                    cell = cell_header(i + 1) + str(index + 1)
                    cell_details = name + ": " + cell
                    sheet_handler(list_details, cell_details, path, sheet_count, row[i])
    except Exception as e:
        file_skip_count += 1
        skipped_file_name.append({"File Name": list_details[0]})
        logging.error("Error has occured while processing this file {} and the error is {}".format(list_details[0], e))
        sys.stdout.write('Error has occurred while processing this file: {}, error is {} \n'.format(list_details[0], e))
        pass


# method which writes all the output data to a list of dictionary - raw_details_list
def sheet_handler(list_details, cell_details, path, page_count, row):
    try:
        global keyword_dic, pattern_list
        if isinstance(row, str):
            for category, list_of_pattern in pattern_list.items():
                for index in range(len(list_of_pattern)):
                    words = list_of_pattern[index].findall(row)
                    for _ in words:
                        temp_dic = {"File Name": list_details[0], "Index Number": list_details[1],
                                    "Folder Name": list_details[2], "File Path": path, "File Type": list_details[3],
                                    "Total Number of Page(s)/Slide(s)": page_count, "Keyword": keyword_dic[category][index],
                                    "IT Category": category, "Page/Slide Number": cell_details}
                        raw_details_list.append(temp_dic)

    except Exception as e:
        logger.info("Error while populating to RAW details sheet {}".format(e))
        sys.stdout.write("Error while populating to RAW details sheet, error is {} \n".format(e))
        sys.exit(1)


# removes unneccessary spaces from extracted texts from images
def removal(text):
    return text.replace("\\r", " ").replace("\\n", " ").replace('\\x', " ").replace('0c', " ")

# handles these extensions - ".PNG", ".png", ".JPEG", ".jpeg", ".JPG", ".jpg", ".gif", ".pnm", ".PNM"
def img_handler(path):
    global file_skip_count
    list_details = folderName(path)
    try:
        text = textract.process(path, encoding='utf8', method='tesseract')
        # text = pytesseract.image_to_string(path, lang='eng')
        content = removal(str(text))
        sheet_handler(list_details, 1, path, 1, content)
    except Exception as e:
        file_skip_count += 1
        skipped_file_name.append({"File Name": list_details[0]})
        logging.error("Error has occurred while processing this file {} and the error is {}".format(list_details[0], e))
        sys.stdout.write('Error has occurred while processing this file: {}, error is {} \n'.format(list_details[0], e))
        pass


# handles these extensions - ".tiff", ".tif", ".jfif", ".bmp"
def tiff_handler(path):
    global file_skip_count
    list_details = folderName(path)
    try:
        out_file = os.path.join(temp_storage_path, "temp_img.JPEG")
        im = Image.open(path)
        out = im.convert("RGB")
        out.save(out_file, "JPEG", quality=90)
        img_handler_for_pdf(out_file, path, 1, 1, list_details)
        os.remove(out_file)
    except Exception as e:
        file_skip_count += 1
        skipped_file_name.append({"File Name": list_details[0]})
        sys.stdout.write('Error has occurred while processing this file: {}, error is {} \n'.format(list_details[0], e))
        logger.info("Error occured while processing {} and error is {}".format(path, e))
        pass


def write_image(shape, count, list_details, total_pages, current_page):
    try:
        image = shape.image
        image_bytes = image.blob
        image_filename = 'image{:03d}.{}'.format(count, image.ext)
        with open(image_filename, 'wb') as fi:
            fi.write(image_bytes)
            fi.close()
        if image_filename.endswith(".wmf"):
            os.remove(image_filename)
        else:
            txt = textract.process(image_filename, encoding='cp437', method='tesseract')
            # txt = pytesseract.image_to_string(Image.open(image_filename))
            sheet_handler(list_details, current_page, path, total_pages, txt)
            os.remove(image_filename)
    except Exception as e:
        sys.stdout.write('Error has occurred while processing this file: {}, error is {} \n'.format(list_details[0], e))
        logger.info("Error occured while processing {} and error is {}".format(path, e))
        pass
        
    
# converts ppt to pdf for further processing
def ppt_handler(path):
    global file_skip_count
    flag_count = 0
    list_details = folderName(path)
    try:
        prs = Presentation(path)
        total_pages = len(prs.slides)
        page = 0
        for slide in prs.slides:
            page = page + 1
            count = 0
            for shape in slide.shapes:
                try:
                    count = count + 1
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        write_image(shape, page, list_details, total_pages, page)
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                # print(run.text)
                                sheet_handler(list_details, page, path, total_pages, run.text)
                except Exception as e:
                    flag_count = flag_count + 1
                    sys.stdout.write('passing this shape:\n'.format(e))
                    logger.info("Error occurred while processing {} and error is shaperror: {}".format(path, e))
                    pass
        print("flsg_count ",flag_count)
        print("total pages ", total_pages)
        if flag_count == total_pages:
            print("inside")
            file_skip_count += 1
            skipped_file_name.append({"File Name": list_details[0]})
    except KeyError as e:
        sys.stdout.write('this is for file rejection:\n{}'.format(e))
        logger.info("Error occurred while processing {} and error is Keyerror: {}".format(path, e))
        file_skip_count += 1
        skipped_file_name.append({"File Name": list_details[0]})
        pass
    except Exception as e:
        sys.stdout.write('Error has occurred while processing this file at file level: {}, error is {} \n'.format(list_details[0], e))
        logger.info("Error occurred while processing {} and error is {}".format(path, e))
        file_skip_count += 1
        skipped_file_name.append({"File Name": list_details[0]})
        pass


# extracts images out of scanned pdfs
def pdf_handler(path):
    global file_skip_count
    list_details = folderName(path)
    try:
        dpi_val = 60
        print("DPI value used:: ", dpi_val)
        pages = convert_from_path(path, dpi=dpi_val)
        total_pages = len(pages)
        for index, page in enumerate(pages):
            temp_file = list_details[0] + str(index + 1) + '.jpeg'
            temp_path = os.path.join(temp_storage_path, temp_file)
            page.save(temp_path, 'JPEG')
            img_handler_for_pdf(temp_path, path, index+1, total_pages, list_details)
            os.remove(temp_path)
    except Exception as e:
        file_skip_count += 1
        skipped_file_name.append({"File Name": list_details[0]})
        sys.stdout.write('Error has occurred while processing this file: {}, error is {} \n'.format(list_details[0], e))
        logger.info("Error occurred while processing {} and error is {}".format(path, e))
        sys.stdout.write("skipping this file\n")
        pass


def docx_handler(path):
    global file_skip_count
    list_details = folderName(path)
    try:
        text = ""
        list_details = folderName(path)
        with open(path, "rb") as docx_file:
            result = mammoth.extract_raw_text(docx_file)
            text += result.value
        sheet_handler(list_details, 1, path, 1, text)
    except Exception as e:
        file_skip_count += 1
        skipped_file_name.append({"File Name": list_details[0]})
        logger.info("exception has occurred while procesing this file:: {} and error is {}".format(path, e))
        sys.stdout.write("exception has occurred while procesing this file:: {}, error is: {}\n".format(path, e))
        print("skipping this file")
        pass

def doc_handler(path):
    global file_skip_count
    list_details = folderName(path)
    letter_count = 0
    flag = False
    try:
        headers = {
            "X-Tika-OCRLanguage": "eng",
            "X-Tika-OCRTimeout": "300"
        }
        if path.endswith((".doc", ".DOC")):
            flag = True
        pages_dict = {}
        data = parser.from_file(path, xmlContent=True, requestOptions={'headers': headers, 'timeout': 300})
        xhtml_data = BeautifulSoup(data['content'], "lxml")
        for i, content in enumerate(xhtml_data.find_all('div', attrs={'class': 'page'})):
            _buffer = StringIO()
            _buffer.write(str(content))
            parsed_content = parser.from_buffer(_buffer.getvalue())

            # Add pages
            text = parsed_content['content']
            # print(text)
            if text:
                pages_dict[i+1] = removal(text.strip())
            else:
                pages_dict[i + 1] = text
            if text:
                letter_count = letter_count + len(removal(text.strip().replace(" ", "")))

        total_pages = len(pages_dict)
        if letter_count >= total_pages*150 or flag:
            for key, value in pages_dict.items():
                sheet_handler(list_details, key, path, total_pages, value)
        else:
            return True
        return False
    except Exception as e:
            file_skip_count += 1
            skipped_file_name.append({"File Name": list_details[0]})
            logger.info("exception has occurred while procesing this file:: {} and error is {}".format(path, e))
            sys.stdout.write("exception has occurred while procesing this file:: {}, error is: {}\n".format(path, e))
            print("skipping this file")
            pass


# handles images converted from pdfs which in turn is converted from docx and ppts
def img_handler_for_pdf(temp_path, path, page_number, total_pages, list_details):
    try:
        # text = textract.process(temp_path, encoding="utf8",  errors='ignore', method='tesseract')
        # text = tools[0].image_to_string(Image.open(temp_path), lang="eng", builder=pyocr.builders.TextBuilder())
        text = pytesseract.image_to_string(temp_path, lang='eng')
        content = removal(str(text))
        sheet_handler(list_details, page_number, path, total_pages, content)
    except Exception as e:
        sys.stdout.write('Error has occurred while processing this file: {}, error is {} \n'.format(list_details[0], e))
        logger.info("Error occurred while processing {} and error is {}".format(path, e))
        sys.stdout.write("skipping this file\n")
        pass


# handles xlsb files
def xlsb_handler(path):
    global file_skip_count
    list_details = folderName(path)
    try:
        sheet_count = len(open_workbook(path).sheets)
        with open_workbook(path) as wb:
            for sheetname in wb.sheets:
                row_number = 0
                with wb.get_sheet(sheetname) as sheet:
                    for row in sheet.rows():
                        row_number = row_number + 1
                        for i in range(len(row)):
                            if row[i].v:
                                cell = cell_header(i + 1) + str(row_number)
                                cell_details = sheetname + ": " + cell
                                sheet_handler(list_details, cell_details, path, sheet_count, row[i].v)
    except Exception as e:
        file_skip_count += 1
        skipped_file_name.append({"File Name": list_details[0]})
        sys.stdout.write('Error has occurred while processing this file: {}, error is {} \n'.format(list_details[0], e))
        logger.info("Error occurred while processing {} and error is {}".format(path, e))
        sys.stdout.write("skipping this file\n")
        pass


# odt file handler
def odt_handler(path):
    global file_skip_count
    list_details = folderName(path)
    try:
        textdoc = load(path)
        allparas = textdoc.getElementsByType(text.P)
        for i in allparas:
            sheet_handler(list_details, 1, path, 1, str(i))
    except Exception as e:
        file_skip_count += 1
        skipped_file_name.append({"File Name": list_details[0]})
        sys.stdout.write('Error has occurred while processing this file: {}, error is {} \n'.format(list_details[0], e))
        logger.info("Error occurred while processing {} and error is {}".format(path, e))
        sys.stdout.write("skipping this file\n")


# processes msg files
def msg_handler(path):
    global file_skip_count
    list_details = folderName(path)
    try:
        content = str(Message(path).body)
        sub = str(Message(path).subject)
        sheet_handler(list_details, 1, path, 1, content)
        sheet_handler(list_details, 1, path, 1, sub)
    except Exception as e:
        file_skip_count += 1
        skipped_file_name.append({"File Name": list_details[0]})
        sys.stdout.write('Error has occurred while processing this file: {}, error is {} \n'.format(list_details[0], e))
        logger.info("Error occurred while processing {} and error is {}".format(path, e))
        sys.stdout.write("skipping this file")


# processes text files
def txt_handler(path):
    list_details = folderName(path)
    global file_skip_count
    try:
        file = open(path, mode='r')
        content = file.read()
        sheet_handler(list_details, 1, path, 1, content)
        file.close()
    except Exception as e:
        file_skip_count += 1
        skipped_file_name.append({"File Name": list_details[0]})
        sys.stdout.write('Error has occurred while processing this file: {}, error is {} \n'.format(list_details[0], e))
        logger.info("Error occurred while processing {} and error is {}".format(path, e))
        sys.stdout.write("skipping this file\n")


# processes rtf files
def rtf_handler(path):
    global file_skip_count
    list_details = folderName(path)
    try:
        file = open(path, mode='r')
        rtf = file.read()
        content = rtf_to_text(rtf)
        sheet_handler(list_details, 1, path, 1, content)
        file.close()
    except Exception as e:
        file_skip_count += 1
        skipped_file_name.append({"File Name": list_details[0]})
        sys.stdout.write('Error has occurred while processing this file: {}, error is {} \n'.format(list_details[0], e))
        logger.info("Error occurred while processing {} and error is {}".format(path, e))
        sys.stdout.write("skipping this file\n")
        pass


# extracts index number from file name
def index_number_verification(index):
    r = re.compile('[@_!#$%^&*()<>?/\|}{~:]')
    if "_" in index:
        index = index.split("_")[0]
    if (index.isupper() or index.islower()) or r.search(index) is not None:
        return " "
    index = re.sub('[^0-9\.\s]+', '', index)
    if len(index) == 1:
        index = re.sub('[.]', " ", index)
    return index


# generates file details
def folderName(path):
    # C:\Users\sb512911\Desktop\All\Screenshots\1.2.1 ILS Structure Chart.csv
    path = os.path.normpath(path)
    file_name = path.split(os.sep)[-1]
    index = ""
    list1 = ['1', '2', '3', '4', '5', '6', '7', '8', '9', '0', '.', '_']
    for i in file_name:
        if i in list1:
            index = index + i
        else:
            break
    index_number = index_number_verification(index)
    folder_name = path.split(file_name)[0]
    file_type = path.split(".")[-1]
    return [file_name, index_number, folder_name, file_type]


# pass the additional parameter for relevance test - dynamically
# generates file level details sheet from "Raw Details"
def generatingFileLevelDetail(file_path):
    try:
        data = pd.read_excel(open(file_path, 'rb'), sheet_name=0)
        if not data.empty:
            data.replace([None, np.nan, "None", "NaN", "nan"], '', inplace=True)
            # print("part -2", data)
            temp_data = data[["File Name", "Index Number", "Folder Name", "File Type"]]
            temp_data["Scan Status"] = "Successful"
            file_level_detail = temp_data[temp_data['File Type'] != "Unknown"].drop_duplicates().sort_values('File Name')
            file_level_detail["Number of Keywords found"] = \
                data[data["File Type"] != "Unknown"][['File Name', "Folder Name"]].sort_values('File Name').groupby(
                    ['File Name', "Folder Name"]).size().reset_index()[0].values
            file_level_detail["Number of Keyword Categories"] = \
                data[data["File Type"] != "Unknown"][['File Name',"Folder Name", 'IT Category']].drop_duplicates().groupby(
                    ['File Name', 'Folder Name']).size().reset_index()[0].values

            intermediate1 = data[data["File Type"] != "Unknown"][['File Name', 'Keyword', "Folder Name"]].groupby(
                ['File Name', "Folder Name", 'Keyword']).size().reset_index()
            intermediate1[[0]] = intermediate1[[0]].astype(str).values
            intermediate1['joined'] = intermediate1[['Keyword', 0]].apply(lambda x: '('.join(x) + "), ", axis=1).values
            file_level_detail["Keyword Details"] = intermediate1[["File Name", "joined", "Folder Name"]].groupby(
                ["File Name", "Folder Name"]).sum().values

            file_level_detail["Category Details"] = \
                data[data["File Type"] != "Unknown"].groupby(["File Name", "Folder Name"]).apply(
                    lambda x: x["IT Category"].drop_duplicates().str.cat(sep=", ")).values
            file_level_detail["File Path"] = data[data["File Type"] != "Unknown"]["File Path"].drop_duplicates().values

            # Relevance check -- code here - assuming relevant count = 5
            mapping = {True: "Yes", False: "No"}

            """
            file_level_detail["Relevant?"] = data[data["File Type"] != "Unknown"][["File Name","Folder Name","Keyword"]]\
                .drop_duplicates()[["File Name", "Folder Name"]].groupby(["File Name", "Folder Name"]).size().\
                reset_index()[0].apply(lambda x: x > 5).map(mapping).values
            """
            file_level_detail["Relevant?"] = \
            data[data["File Type"] != "Unknown"][["File Name", "Folder Name", "Keyword"]].groupby(["File Name", "Folder Name"]).size(). \
                reset_index()[0].apply(lambda x: x >= 5).map(mapping).values


            temp_data_unknown = temp_data[["File Name", "Index Number", "Folder Name", "File Type"]][
                temp_data['File Type'] == "Unknown"]
            temp_dict = {"Number of Keywords found": 0, "Number of Keyword Categories": 0, "Keyword Details": "",
                         "Category Details": "", "Relevant?": "", "File Path": "", "Scan Status": "Unsuccessful"}
            temp_df = pd.DataFrame([temp_dict])
            # print("part4", temp_df)
            unknown_files = cross_join(temp_data_unknown, temp_df)
            file_level_detail = file_level_detail.append(unknown_files, ignore_index=True)

            #file_level_detail = file_level_detail.replace(np.nan, '', regex=True)
            #file_level_detail.fillna('', inplace=True)

            file_level_detail.replace([None, np.nan, "None", "NaN", "nan"], '', inplace=True)
            sys.stdout.write("Total processed Files including unknown extension:: {}\n".format(len(file_level_detail)))
            writing_to_excel(file_level_detail, "File Level Detail", "a", file_path)

    except Exception as e:
        sys.stdout.write("error occurred while generating file level detail - Raw details doesnt have any relevant value \n")
        sys.stdout.write("And the error is: {} \n".format(e))
        logger.info("error occurred while generating file level detail, error is {}".format(e))
        # sys.exit(1)


# generates keyword detail output sheet sheet
def generatingKeywordDetails(file_path):
    try:
        data = pd.read_excel(open(file_path, 'rb'), sheet_name="Raw Details")
        data = data[data['File Type'] != "Unknown"]

        data1 = data[["Keyword"]]
        data1.replace([None, np.nan, "None", "NaN", "nan", ''], np.nan, inplace=True)
        if len(data1.dropna()) != 0:
            data.replace([None, np.nan, "None", "NaN", "nan"], '', inplace=True)

            # distinct keywords and categories
            keyword_details = data[["Keyword", "IT Category"]].drop_duplicates()
            # keyword_details.replace([None, np.nan, "None", "NaN", "nan"], '', inplace=True)
            # print(keyword_details)

            #  Number of files with keywords
            keyword_details = data[["Keyword", "File Name"]].drop_duplicates().groupby(["Keyword"]) \
                .size().reset_index().merge(keyword_details, on="Keyword")

            # renaming
            keyword_details = keyword_details. \
                rename(columns={'IT Category': 'Keyword Category', 0: 'Number of files with keywords'})

            # Total keyword hits
            keyword_details = data[["File Name", "Keyword"]].groupby(["Keyword"]) \
                .size().reset_index().merge(keyword_details, on="Keyword").rename(columns={0: "Total keyword hits"})

            # File(s) with maximum keyword hits and File Path(s)
            intermediate2 = data[["Keyword", "File Name", "File Path"]].groupby(["Keyword", "File Name", "File Path"]) \
                .size().reset_index().merge(keyword_details, on="Keyword") \
                .rename({0: "Count"}, axis=1)[["File Name", "Keyword", "File Path", "Count"]]

            intermediate3 = intermediate2.groupby('Keyword')['Count'].apply(lambda x: x.eq(x.max()))

            intermediate4 = intermediate2.loc[intermediate3].groupby(['Keyword'])['File Name'].agg(
                ', '.join).reset_index()

            intermediate5 = intermediate2.loc[intermediate3].groupby(['Keyword'])['File Path'].agg(
                ', '.join).reset_index()

            keyword_details = keyword_details.merge(intermediate4, on="Keyword")

            keyword_details = keyword_details.merge(intermediate5, on="Keyword")

            keyword_details = keyword_details.rename(columns={"File Path": "File Path(s) with maximum keyword",
                                                              "File(s) with maximum keyword hits":
                                                                  "File Path(s) with maximum keyword"})

            # All file(s) with keyword hits
            keyword_details["File Path(s)"] = data[["Keyword", "File Path"]].groupby(["Keyword"]) \
                .apply(lambda x: x["File Path"].drop_duplicates().str.cat(sep=", ")).reset_index() \
                .merge(keyword_details, on="Keyword")[[0]].values

            # Index Number(s)
            keyword_details["Index Number(s)"] = \
                data[["Keyword", "Index Number"]].groupby("Keyword") \
                    .apply(lambda x: x["Index Number"].drop_duplicates().astype(str).str.cat(sep=", ")).reset_index() \
                    .merge(keyword_details, on="Keyword")[[0]].values

            # All the file names which has the keyword
            keyword_details["File Name(s)"] = data[["Keyword", "File Name"]] \
                .groupby("Keyword").apply(lambda x: x["File Name"].drop_duplicates()
                                          .str.cat(sep=", ")).reset_index().merge(keyword_details, on="Keyword")[
                [0]].values

            keyword_details = keyword_details.rename(columns={"File Name": "File Name(s) with maximum keyword",
                                                              "File Name(s)": "All file(s) with keyword hits"})

            keyword_details = keyword_details[["Keyword", "Keyword Category", "Number of files with keywords",
                                               "Total keyword hits", "File Name(s) with maximum keyword",
                                               "File Path(s) with maximum keyword",
                                               "All file(s) with keyword hits", "Index Number(s)", "File Path(s)"]]

            # writing to "File Level Detail" sheet
            # keyword_details = keyword_details.replace(np.nan, '', regex=True)
            # keyword_details.fillna('', inplace=True)
            keyword_details.replace([None, np.nan, "None", "NaN", "nan"], '', inplace=True)
            writing_to_excel(keyword_details, "Keyword Details", "a", file_path)

    except Exception as e:
        sys.stdout.write("Error occured while generating Keyword details - keyword level details could not be generated\n")
        sys.stdout.write("And the error is: {} \n".format(e))
        logger.info("Error occured while generating Keyword details - keyword level details could not be generated"
                    "and the error is: {}".format(e))
        pass


# handles DRL file and writes to DRL Details sheet
def generatingDRLDetails(drl_path, output_path):
    global warning_message
    df = pd.read_excel(output_path, sheet_name="Raw Details")
    try:
        if not df.empty:
            df.replace([None, np.nan, "None", "NaN", "nan"], '', inplace=True)
            data = df[df["File Type"] != "Unknown"][["File Name", "Index Number", "File Path", "Keyword"]].drop_duplicates()
            dict_file_path = {}
            all_files = data["File Name"].tolist()
            dict_all = {}
            dict_index_number = {}
            each_row = []
            for file in all_files:
                dict_all[file] = data[data['File Name'] == file]['Keyword'].tolist()
                dict_index_number[file] = data[data['File Name'] == file]['Index Number'].tolist()[0]
                dict_file_path[file] = data[data['File Name'] == file]['File Path'].tolist()[0]
            workbook1 = xlrd.open_workbook(drl_path)
            sh = workbook1.sheet_by_index(0)
            for rownum in range(sh.nrows):
                rows = sh.row_values(rownum)
                if rows[5].strip().lower() == "it":
                    sentence = rows[6]
                    drl_number = int(rows[1])
                    satisfied_output_file_name = {}
                    satisfied_output_file_path = {}
                    satisfied_output_index = {}
                    flag = "No"
                    relevant = 0
                    # keys = file name
                    for keys, values in dict_all.items():
                        count = 0
                        for m in values:
                            r1 = re.compile('|'.join([r'(?<!\w)%s(?!\w)' % re.escape(m)]), flags=re.I)
                            matches = r1.findall(sentence)
                            if matches:
                                count = count + 1
                        if count >= 5:
                            flag = "Yes"
                            relevant = relevant + 1
                            if drl_number not in satisfied_output_file_name:
                                satisfied_output_file_name[int(drl_number)] = keys + ', '
                                satisfied_output_file_path[int(drl_number)] = dict_file_path[keys] + ", "
                                satisfied_output_index[int(drl_number)] = str(dict_index_number[keys]) + ", "
                            else:
                                satisfied_output_file_name[int(drl_number)] += keys + ", "
                                satisfied_output_file_path[int(drl_number)] += dict_file_path[keys] + ", "
                                satisfied_output_index[int(drl_number)] += str(dict_index_number[keys]) + ", "
                    dict1 = {}
                    dict1["DRL #"] = drl_number
                    dict1["Request"] = sentence
                    dict1["Relevant Files Found?"] = flag
                    dict1["Number of relevant files"] = relevant
                    if satisfied_output_file_name and satisfied_output_file_path and satisfied_output_index:
                        dict1["File Name(s)"] = satisfied_output_file_name[drl_number]
                        dict1["Index Number(s)"] = satisfied_output_index[drl_number]
                        dict1["File Path(s)"] = satisfied_output_file_path[drl_number]
                        each_row.append(dict1)
                    else:
                        dict1["Index Number(s)"] = ""
                        dict1["File Path(s)"] = ""
                        dict1["File Name(s)"] = ""
                        each_row.append(dict1)

            dataframe = pd.DataFrame(each_row)
            dataframe = dataframe[
                    ["DRL #", "Request", "Relevant Files Found?", "Number of relevant files", "File Name(s)", "Index Number(s)",
                     "File Path(s)"]]
            # writing to "File Level Detail" sheet
            #dataframe = dataframe.replace(np.nan, '', regex=True)
            #dataframe = dataframe.fillna('', inplace=True)
            dataframe.replace([None, np.nan, "None", "NaN", "nan"], '', inplace=True)
            writing_to_excel(dataframe, "DRL Details", "a", output_path)
    except Exception as e:
        warning_message = warning_message + "|" + "DRL File structure is incorrect"
        sys.stdout.write("Error has occured while generating DRL file sheet, error is {} \n".format(e))
        sys.stdout.write('DRL file structure is incorrect \n')
        logger.info("Please check your DRL file structure and the specific error is {}".format(e))
        pass
        #sys.exit(1)


# writes, formats, colours the final output file
def writing_to_excel(dataframe, sheet_name, mode, output_file):
    try:
        if not os.path.exists(output_file):
            open(output_file, 'w').close()
        with pd.ExcelWriter(output_file, engine='openpyxl', mode=mode) as writer:
            dataframe.to_excel(writer, sheet_name, index=False)
            writer.save()

        wb1 = openpyxl.load_workbook(output_file)
        ws = wb1[sheet_name]
        fillBack = PatternFill(start_color="00FFFF00", fill_type="solid")

        ws.column_dimensions['A'].width = 24.94
        ws.column_dimensions['B'].width = 24.94
        ws.column_dimensions['C'].width = 24.94
        ws.column_dimensions['D'].width = 24.94
        ws.column_dimensions['E'].width = 24.94
        ws.column_dimensions['F'].width = 24.94
        ws.column_dimensions['G'].width = 24.94
        ws.row_dimensions[1].height = 28.8

        for cell in ws["1:1"]:
            cell.fill = fillBack
        wb1.save(output_file)

    except Exception as e:
        logger.info("Exception occurred while generating this sheet:: {}".format(sheet_name))
        print("Exception occurred while generating this sheet:: {}".format(sheet_name))
        pass


# does a cross join in python
def cross_join(left, right):
    return left.assign(key=1).merge(right.assign(key=1), on='key').drop('key', 1)


# creates regex pattern for each keyword
def pattern(keys):
    return re.compile('|'.join([r'(?<!\w)%s(?!\w)' % re.escape(keys)]), flags=re.I)


# creating a dictionary of lists containing the regex pattern for each keyword
def keywordlist(team):
    global keyword_dic, pattern_list
    df = keyword_extraction(team)
    category_list = df["category"].drop_duplicates().tolist()
    for i in category_list:
        temp = df[df["category"] == i]["keyword"].tolist()
        keyword_dic[i] = temp
        pattern_list[i] = [pattern(key) for key in temp]


# extracts the data dictionary from mysql table - data_dictionary.virtual_data_room
def keyword_extraction(team):
    try:
        # db_connection_str = 'mysql+pymysql://root:admin@localhost/virtual_data_room'
        db_connection_str = 'mysql+pymysql://' + user + ":" + password + "@" + address + "/" + database_name
        db_connection = create_engine(db_connection_str)
        query = 'SELECT keyword, category FROM {} where team = "{}"'.format(table_name, team)
        df = pd.read_sql(query, con=db_connection)
        #print(df)
        return df.drop_duplicates(['keyword'])
    except Exception as e:
        logger.info("mysql connection issue: data dictionary could not be accessed, error is: {}".format(e))
        sys.stdout.write("mysql connection issue: data dictionary could not be accessed, error is: {} \n".format(e))
        sys.exit(1)


def multi_processing(i):
    global file_skip_count
    sys.stdout.write("file name {} \n".format(i))
    logger.info("processing file:: {}".format(i))
    if i.endswith(".csv"):
        csv_handler(i)
    elif i.endswith(".xlsb"):
        xlsb_handler(i)
    elif i.endswith(".msg"):
        msg_handler(i)
    elif i.endswith(".txt"):
        txt_handler(i)
    elif i.endswith(".rtf"):
        rtf_handler(i)
    elif i.endswith((".docx", '.DOCX')):
        sys.stdout.write("{}\n".format(datetime.datetime.now()))
        docx_handler(i)
    elif i.endswith((".doc", ".DOC")):
        doc_handler(i)
    elif i.endswith(".odt"):
        odt_handler(i)
    elif i.endswith((".pptx", "ppt", ".PPTX", ".PPT")):
        sys.stdout.write("{}\n".format(datetime.datetime.now()))
        ppt_handler(i)
    elif i.endswith((".PNG", ".png", ".JPEG", ".jpeg", ".JPG", ".jpg", ".gif", ".pnm", ".PNM")):
        img_handler(i)
    elif i.endswith((".xlsx", ".xlsm", ".xls", ".XLS", "XLSX")):
        xlsx_handler(i)
    elif i.endswith((".pdf", ".PDF")):
        sys.stdout.write("{}\n".format(datetime.datetime.now()))
        flag = doc_handler(i)
        if flag:
            pdf_handler(i)
    elif i.endswith((".tiff", ".tif", ".jfif", ".bmp")):
        tiff_handler(i)
    else:
        # files which are having extension other than the above ones, will be declared Unknown
        sys.stdout.write("File {} is not handled by this application\n".format(path))
        file_skip_count = file_skip_count + 1
        list_details = folderName(i)
        temp_dic = {"File Name": list_details[0], "Index Number": list_details[1],
                    "Folder Name": list_details[2], "File Path": i, "File Type": "Unknown",
                    "Total Number of Page(s)/Slide(s)": "", "Keyword": "", "IT Category":
                    "", "Page/Slide Number": ""}
        raw_details_list.append(temp_dic)
        skipped_file_name.append({"File Name": list_details[0]})


def file_details(s3_directory, guid, user_folder, folder):
    bucket_name = "aws-workdocs-test-bucket"
    try:
        session = Session(aws_access_key_id=access_key,
                          aws_secret_access_key=secret_key)
        s3 = session.resource('s3')
        your_bucket = s3.Bucket(bucket_name)
        if not os.path.exists(user_folder):
            os.makedirs(user_folder, exist_ok=True)
        for s3_file in your_bucket.objects.all():
            if s3_file.key.startswith(os.path.join(guid, folder)):
                filename = s3_file.key.split("/")[-1]
                if not filename:
                    continue
                path = s3_file.key.split(filename)[0][:-1]
                path_local = os.path.join(user_folder, path)
                if not os.path.exists(path_local):
                    os.makedirs(path_local)
                os.chdir(path_local)
                logger.info("file name to be downloaded {}".format(filename))
                logger.info("file path in local {}".format(path_local))
                s3.meta.client.download_file(bucket_name, s3_file.key, filename)
        return os.path.join(user_folder, guid, folder)

    except Exception as e:
        logger.info("Error has occurred while downloading the files from S3 directory, error is {}".format(e))
        sys.stdout.write("Error has occurred while downloading the files from S3 directory, error is {}\n".format(e))
        sys.exit(1)


def upload_to_s3(output_file, bucket, object_name):
    session = Session(aws_access_key_id=access_key,
                      aws_secret_access_key=secret_key)
    s3 = session.resource('s3')
    try:
        sys.stdout.write("output upload directory: {} \n".format(object_name))
        s3.meta.client.upload_file(output_file, bucket, object_name)
    except Exception as e:
        sys.stdout.write("Uploading to S3 has caused an issue, errro is {} \n".format(e))
        sys.stdout.write("Output file could not be shared with the application \n")
        return False
    return True


def guid_validation(guid, vdr, drl):
    guid_vdr = vdr.split("/")[-2]
    guid_drl = drl.split("/")[-2]
    if guid == guid_drl == guid_vdr:
        pass
    else:
        logger.info("guids are not matching")
        sys.stdout.write("guids are not matching\n")
        sys.exit(1)

def get_team_name(guid):
    try:
        # db_connection_str = 'mysql+pymysql://root:admin@localhost/virtual_data_room'
        db_connection_str = 'mysql+pymysql://' + user + ":" + password + "@" + address + "/" + database_name
        db_connection = create_engine(db_connection_str)
        query = 'SELECT team FROM user_details where GUID = "{}"'.format(guid)
        df = pd.read_sql(query, con=db_connection)
        return df["team"].tolist()[0]

    except Exception as e:
        logger.info("mysql connection issue: data dictionary could not be accessed: {}".format(e))
        sys.stdout.write("mysql connection issue: data dictionary could not be accessed\n")
        sys.exit(1)



if __name__ == '__main__':
    # user input acceptance
    print('newer code has been updated with 150 chars per page')
    root = os.environ["S3_VDR"]  # s3://aws-workdocs-test-bucket/512911/VDR_345
    sys.stdout.write("Root: {} \n".format(root))
    drl_folder = os.environ["S3_DRL"]  # s3://aws-workdocs-test-bucket/512911/DRL_543
    sys.stdout.write("DRL: {} \n".format(drl_folder))
    keywords = os.environ["ADD_KEYS"]  # TU, Oracle
    sys.stdout.write("keywords: {}\n".format(keywords))
    guid = os.environ["GUID"]  # guid
    sys.stdout.write("GUID: {}\n".format(guid))
    email_list = os.environ["EMAIL_RCPT"]
    sys.stdout.write("Email list: {}\n".format(email_list))
    workdocsDir = os.environ["VDR_DIR_NM"]
    sys.stdout.write("New Param: {}\n".format(workdocsDir))

    try:
        guid_validation(guid, root, drl_folder)
        # VDR folder extraction
        folder = root.split("/")[-1]
        vdr = file_details(root, guid, user_folder, folder)
        folder = drl_folder.split("/")[-1]
        fileList = []
        for path, subdirs, files in os.walk(vdr):
            for name in files:
                a = os.path.join(path, name)
                fileList.append(a)
        logger.info("All files in docker local to be processed are:: {}".format(fileList))
        if len(fileList) == 0:
            sys.stdout.write("VDR directory is empty\n")
            logger.info("VDR directory is empty")
            sys.exit(1)


        # DRL folder details extraction
        drl_file_path = ""
        drl_folder = file_details(drl_folder, guid, user_folder, folder)
        for path, subdirs, files in os.walk(drl_folder):
            for name in files:
                a = os.path.join(path, name)
                drl_file_path = a
        flag = True
        if not drl_file_path:
            sys.stdout.write("DRL directory is empty\n")
            flag = False
            logger.info("DRL directory is empty")
            warning_message= warning_message + "|" + "DRL directory is empty"
            pass


        # data dictionary extraction
        team_name = get_team_name(guid)
        keywordlist(team_name)
        print("additional Keywords ", keywords)
        a = ["None", "none", "NaN", "nan", "NAN"]
        if keywords.strip():
            if keywords.lower() not in a:
                additional_keywords = [i.strip() for i in keywords.split(",")]
                keyword_dic["User Input"] = additional_keywords
                pattern_list["User Input"] = [pattern(key) for key in additional_keywords]

        # processing the files
        for file in fileList:
            multi_processing(file)

        # generating output directory and file in local
        output_directory = os.path.join(user_folder, guid, "output")
        if not os.path.exists(output_directory):
            os.makedirs(output_directory, exist_ok=True)
        output_file = os.path.join(output_directory, "output.xlsx")

        # removing the S3 bucket name from the file path
        path_break = os.path.join(user_folder, guid)
        for index, items in enumerate(raw_details_list):
            items["Folder Name"] = items["Folder Name"].split(path_break)[1]
            items["File Path"] = items["File Path"].split(path_break)[1]
            raw_details_list[index] = items

        # print("raw details sheet generation")
        # generating and formatting the raw details page
        raw_dataframe = pd.DataFrame(raw_details_list)
        print("Total Keyword hit count :: ", len(raw_dataframe))
        writing_to_excel(raw_dataframe, "Raw Details", "o", output_file)

        #print("before generating other sheets")
        # generating the other three output sheets, 2 in case DRL file is absent/corrupt
        generatingFileLevelDetail(output_file)
        generatingKeywordDetails(output_file)
        if flag:
            generatingDRLDetails(drl_file_path, output_file)

        # files with no keywords found
        skipped_file_list = [i["File Name"] for i in skipped_file_name]
        for i in range(len(fileList)):
            temp = os.path.normpath(fileList[i])
            file_name = temp.split(os.sep)[-1]
            all_files.append(file_name)

        data = pd.read_excel(open(output_file, 'rb'), sheet_name=0)
        files_with_matching_keywords = []
        if not data.empty:
            files_with_matching_keywords = data["File Name"].drop_duplicates().tolist()

        additional_files = files_with_matching_keywords + skipped_file_list
        list_of_files_with_no_keywords = list(set(all_files) - set(additional_files))
        list_df = pd.DataFrame({"File Name": list_of_files_with_no_keywords})
        sys.stdout.write("Files with no matching keywords count is:: {}\n".format(len(list_df)))
        writing_to_excel(list_df, "Files with no matching keywords", "a", output_file)

        # unknown extension + file skipped
        unknown_extension_sheet = pd.DataFrame(skipped_file_name)
        sys.stdout.write("Skipped file count is:: {} \n".format(len(unknown_extension_sheet)))
        writing_to_excel(unknown_extension_sheet, "Skipped Files", "a", output_file)

        # output formatting for sending to API
        s3_path = root.split(guid)[0][:-1]
        out = os.path.join(guid, "output", "{}_output.xlsx".format(guid))
        bucket_name = root.split("//")[1].split("/")[0]


        # uploading the output file in S3 and api call to front end
        if upload_to_s3(output_file, bucket_name, out):
            out = os.path.join(s3_path, guid, "output", "{}_output.xlsx".format(guid))
            data_set = {"url": out, "guid": guid, "number_of_file_skipped": file_skip_count, "message": warning_message
                        , "email_list": email_list, "workdocsDir": workdocsDir}
            json_dump = json.dumps(data_set)
            url = 'https://itduediligence.pwcinternal.com/data-room-sweeper/dataRoom/ProcessEmailRequest/'
            headers = {'content-type': 'application/json', 'Accept-Charset': 'UTF-8'}
            result = requests.post(url, data=json_dump, headers=headers, verify=False)
            if result.status_code == HTTPStatus.OK:
                pass
            else:
                sys.stdout.write("Error occurred while downloading the output file or sending the email\n")
                sys.exit(1)
                       
        else:
            sys.stdout.write("Output file could not be shared with the application\n")
            logger.info("Output file could not be shared with the application")
            sys.exit(1)

        # print("path", os.path.join(user_folder, guid))
        shutil.rmtree(os.path.join(user_folder, guid))
        

    except KeyError as e:
        sys.stdout.write("Either data is not present or column is missing and the error is: {} \n".format(e))
        logger.info("Either data is not present or column is missing:: ".format(e))
        sys.exit(1)
    except PermissionError as e:
        sys.stdout.write("Some files are open not have permission to open: {} \n".format(e))
        logger.info("Some files are open not have permission to open:: {}".format(e))
        sys.exit(1)
    except Exception as e:
        logger.info("Fatal error has occured:: ".format(e))
        sys.stdout.write("Error has occured, error is: {} \n".format(e))
        sys.exit(1)
    finally:
        # uploading the log file to S3
        logger.info("Execution finished time:: {}".format(datetime.datetime.now()))
        logger.info("The end")
        logger.info("***********************************NEXT RUN************************************")
        logger.info("*******************************************************************************")
        logger.info("*******************************************************************************")


        bucket_name = "aws-workdocs-test-bucket"
        out = os.path.join("Logs", guid, "{}".format(timestr))
        try:
            upload_to_s3(log_file, bucket_name, out)
        except Exception as e:
            sys.stdout.write("Error has occurred while uploading the logs to S3 and the error is:: {}\n".format(e))
        path = os.path.join("s3://aws-workdocs-test-bucket", out)
        sys.stdout.write("Processing log can be found in S3 in this location:: {}\n".format(path))
        sys.stdout.write("End time: {}".format(datetime.datetime.now()))

