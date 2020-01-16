# needs enhancement - but works fine 
# imagemagick and tesseract-ocr has to be installed separately

import os
import tempfile
import pdf2image
import ctypes
import wand.image
from wand.api import library
import gc
import sys
import re
from pptx import Presentation
import docxpy
from pptx.enum.shapes import MSO_SHAPE_TYPE
import glob
import pandas as pd
from pandas import ExcelWriter
from pandas import ExcelFile
import xlsxwriter
import datetime
from docx import Document
from PIL import Image as PILImage
import pyocr
import pyocr.builders
from zipfile import ZipFile
import shutil
import pytesseract 
from pdf2image import convert_from_path 

now = str(datetime.datetime.now())

tools = pyocr.get_available_tools()

library.MagickNextImage.argtypes = [ctypes.c_void_p]
library.MagickNextImage.restype = ctypes.c_int

if len(tools) == 0:
    raise Exception("No OCR tool found")
tool = tools[0]
langs = tool.get_available_languages()
lang = langs[0]

walk_dir = sys.argv[1]

id45 = ['General Ledger','Accounts Payable','Accounts Receivable','Invoicing','Treasury','CRM','Customer Relationship','Purchasing','Payroll','Attendance','Talent Acquisition','Recruiting','Hiring','Learning and Development','Learning & Development','Document Management','Dispatch','Project Management','Job Management','Scheduling']
id46 = ['IT project','IT initiative','IT budget']
id47 = ['IT expansion','IT investment','business plan']
id137 = ['LAN','WAN','network infrastructure','data center','data centre','operating system','messaging system','file server','print serv','high availability','business continuity','disaster recovery','BCP','DRP']

r = re.compile('|'.join([r'\b%s\b' % w for w in id45]), flags=re.I)
r2 = re.compile('|'.join([r'\b%s\b' % w for w in id46]), flags=re.I)
r3 = re.compile('|'.join([r'\b%s\b' % w for w in id47]), flags=re.I)
r4 = re.compile('|'.join([r'\b%s\b' % w for w in id137]), flags=re.I)

workbook = xlsxwriter.Workbook('VDR_Index.xlsx')
worksheet = workbook.add_worksheet()
worksheet.write('A1','Project Name Here: Document Index')
worksheet.write('A2','As of ' + str(now))
worksheet.write('A3','Kind')
worksheet.write('B3','#')
worksheet.write('C3','ID')
worksheet.write('D3','Document Size')
worksheet.write('E3','Title')
worksheet.write('F3','Folder Path')
worksheet.write('G3','Reviewed/Not Reviewed')
worksheet.write('H3','Status')
worksheet.write('I3','In Scope / Out of Scope')
worksheet.write('J3','PwC Notes')
worksheet.write('K3','Protection')
worksheet.write('L3','Document Link')

n=0
c=3


# Handling/finding in these files : jpg, gif, 

def write_image(shape, text_filename):
    global n
    image = shape.image
    # ---get image "file" contents---
    image_bytes = image.blob
    # ---make up a name for the file, e.g. 'image.jpg'---
    image_filename = 'image{:03d}.{}'.format(n, image.ext)
    n += 1
    with open(image_filename, 'wb') as fi:
        fi.write(image_bytes)
        fi.close()
    if image_filename.endswith(".gif"):
        os.system('convert ' + image_filename + ' ' + 'image.jpg')
        txt = pytesseract.image_to_string(PILImage.open('image.jpg'))
        os.remove('image.jpg')
        os.remove(image_filename)
    elif image_filename.endswith(".wmf"):
        os.system('unoconv  -f pdf -o image.pdf ' + image_filename)
        os.system('convert -density 300 -trim -bordercolor white -border 5 image.pdf image.png')
        txt = pytesseract.image_to_string(PILImage.open('image.png'))
        os.remove('image.png')
        os.remove('image.pdf')
        os.remove(image_filename)
    else: 
        txt = pytesseract.image_to_string(PILImage.open(image_filename))
        os.remove(image_filename)
    searchimgtext(txt, text_filename)


	
def match_search(content,fullmatches,unique_matches):
    matches = r.search(content)
    if matches:
        fullmatches.append(matches.group())
    matches2 = r2.search(content)
    if matches2:
        fullmatches.append(matches2.group())
    matches3 = r3.search(content)
    if matches3:
        fullmatches.append(matches3.group())
    matches4 = r4.search(content)
    if matches4:
        fullmatches.append(matches4.group())

def searchimgtext(content,text_filename):
    fullmatches = []
    unique_matches = []
    match_search(content,fullmatches,unique_matches)
    if fullmatches:
        matchset = set(fullmatches)
        unique_matches = (list(matchset))
        text_filename.write("From images: " + str(unique_matches))

def write_image2(image_filename,c,file):
    #global n
    #image = shape.image
    # ---get image "file" contents---
    #image_bytes = image.blob
    # ---make up a name for the file, e.g. 'image.jpg'---
    #image_filename = 'image{:03d}.{}'.format(n, image.ext)
    #n += 1
    if image_filename.endswith(".gif"):
       os.system('convert ' + image_filename + ' ' + image_filename+'.jpg')
       image_filename = image_filename + '.jpg'
    txt = tool.image_to_string(PILImage.open(image_filename),lang="eng",builder=pyocr.builders.TextBuilder())
    #print(txt)
    searchimgtext2(txt,c,filetype)
    #Image.close(image_filename)

def searchimgtext2(content,c,filetype):
    fullmatches = []
    unique_matches = []
    match_search(content,fullmatches,unique_matches)
    if fullmatches:
        match_report(filetype,filename,file_path,c,fullmatches)
    else:
        no_match_report(filetype,filename,file_path,c)

def write_text(shape):
    text_filename.write(shape.text)

def visitor(shape,text_filename):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            visitor(s,text_filename)
    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        write_image(shape,text_filename)
    elif hasattr(shape, "text"):
        write_text(shape)

def searchtext(file_path,c,filetype):
     with open(file_path, encoding = 'ISO-8859-1') as infile:
         fullmatches = []
         unique_matches = []
         for line in infile.readlines():
             matches = r.search(line)
             if matches:
                 fullmatches.append(matches.group())
             matches2 = r2.search(line)
             if matches2:
                 fullmatches.append(matches2.group())
             matches3 = r3.search(line)
             if matches3:
                 fullmatches.append(matches3.group())
             matches4 = r4.search(line)
             if matches4:
                 fullmatches.append(matches4.group())
         if fullmatches:
             match_report(filetype,filename,file_path,c,fullmatches)
         else:
             no_match_report(filetype,filename,file_path,c)
         infile.close()
         os.remove(file_path)

def searchtext2(content,c,filetype,filename,file_path,fullmatches,unique_matches):
    #fullmatches = []
    #unique_matches = []
    match_search(content,fullmatches,unique_matches)
    if fullmatches:
        match_report(filetype,filename,file_path,c,fullmatches)
    else:
        no_match_report(filetype,filename,file_path,c)

def searchtext3(content,text_filename):
    fullmatches = []
    unique_matches = []
    match_search(content,fullmatches,unique_matches)
    if fullmatches:
        text_filename.write(fullmatches)

def searchpdf(content,fullmatches,unique_matches):
     matches = r.search(content)
     if matches:
        fullmatches.append(matches.group())
     matches2 = r2.search(content)
     if matches2:
         fullmatches.append(matches2.group())
     matches3 = r3.search(content)
     if matches3:
         fullmatches.append(matches3.group())
     matches4 = r4.search(content)
     if matches4:
         fullmatches.append(matches4.group())

def searchdocx(file_path,c,filename):
     fullmatches = []
     unique_matches = []
     filetype = "Microsoft Word"
     doc = Document(file_path)
     for i in doc.paragraphs:
         matches = r.search(i.text)
         if matches:
             fullmatches.append(matches.group())
         matches2 = r2.search(i.text)
         if matches2:
             fullmatches.append(matches2.group())
         matches3 = r3.search(i.text)
         if matches3:
             fullmatches.append(matches3.group())
         matches4 = r4.search(i.text)
         if matches4:
             fullmatches.append(matches4.group())
     if fullmatches:
         match_report(filetype,filename,file_path,c,fullmatches)
     else:
         no_match_report(filetype,filename,file_path,c)

def no_match_report(filetype,filename,file_path,c):
    worksheet.write('A'+str(c),filetype)
    worksheet.write('E'+str(c),filename)
    worksheet.write('F'+str(c),file_path)
    worksheet.write('G'+str(c),'Scanned')
    worksheet.write('H'+str(c),'On File System')
    worksheet.write('I'+str(c),'Out of Scope')

def match_report(filetype,filename,file_path,c,fullmatches):
    matchset = set(fullmatches)
    unique_matches = (list(matchset))
    worksheet.write('A'+str(c),filetype)
    worksheet.write('E'+str(c),filename)
    worksheet.write('F'+str(c),file_path)
    worksheet.write('G'+str(c),'Scanned')
    worksheet.write('H'+str(c),'On File System')
    worksheet.write('I'+str(c),'In Scope')
    worksheet.write('J'+str(c),'Matches: ' + str(unique_matches))

for root, subdirs, files in os.walk(walk_dir):
    list_file_path = os.path.join(root, 'my-directory-list.txt')

    with open(list_file_path, 'wb') as list_file:
        for subdir in subdirs:
            c += 1
            worksheet.write('A'+str(c),'Folder')
            worksheet.write('E'+str(c),subdir)
            worksheet.write('F'+str(c),root+'\\'+subdir)
            print(str(c)+": "+subdir)
        for filename in files:
            gc.collect()
            c += 1
            print(str(c)+": "+filename)
            if filename.endswith((".xlsx")):
                 fullmatches = []
                 unique_matches = []
                 content = []
                 filetype = 'Microsoft Excel'
                 file_path = os.path.join(root, filename)
                 df = pd.read_excel(file_path)
                 for column in df.columns:
                     for i in df.index:
                         content.append(df[column][i])
                 searchtext2(str(content),c,filetype,filename,file_path,fullmatches,unique_matches)
            elif filename.endswith(".doc"):
                 file_path = os.path.join(root, filename)
                 file_path = file_path.replace(' ','\ ')
                 docx_file = file_path + 'x' 
                 os.system('antiword ' + '\"'+file_path+'\"' + ' > ' + '\"'+docx_file+'\"')
                 searchdocx(docx_file,c,filename)
                 os.remove(docx_file)
                 #word = win32.gencache.EnsureDispatch('Word.Application')
                 #doc = word.Documents.Open(file_path)
                 #doc.Activate()
                 #new_file_abs = os.path.abspath(file_path)
                 #new_file_abs = re.sub(r'\.\w+$', '.docx', new_file_abs)
                 #word.ActiveDocument.SaveAs(
                 #    new_file_abs, FileFormat=constants.wdFormatXMLDocument
                 #)
                 #doc.Close(False)
                 #os.remove(file_path)
                 #searchdocx(new_file_abs,c,filename)
            elif filename.endswith((".docx",".DOCX")):
                 file_path = os.path.join(root, filename)
                 searchdocx(file_path,c,filename)
            elif filename.endswith((".txt")) and not filename.endswith("my-directory-list.txt"):
                 filetype = 'Text File'
                 file_path = os.path.join(root, filename)
                 searchtext(file_path,c,filetype)
            elif filename.endswith((".jpg",".png",".bmp",".pnm",".jfif",".jpeg",".tiff")):
                 filetype = 'Image'
                 file_path = os.path.join(root, filename)
                 write_image2(file_path,c,filetype)
            elif filename.endswith((".pptx", ".ppt")):
                 filetype = 'Microsoft Powerpoint'
                 file_path = os.path.join(root, filename)
                 text_filename = open(file_path+'_raw.txt', 'w')
                 prs = Presentation(file_path)
                 for slide in prs.slides:
                     for shape in slide.shapes:
                         if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                             for s in shape.shapes:
                                 visitor(s,text_filename)
                         if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                             write_image(shape,text_filename)
                         if shape.has_text_frame:
                             write_text(shape)
                 with ZipFile(file_path) as inzip:
                    embed_list = [name for name in inzip.namelist()]
                    inzip.extractall(root+'/tempout', embed_list)
                    for temproot, tempsubdirs, tempfiles in os.walk(root+'ppt\\embeddings'):
                        for tempfile in tempfiles:
                            tempcontent = []
                            tempname = info.tempfile
                            #print(tempname)
                            if tempname.endswith((".xls", ".xlsx")):
                                df = pd.read_excel(tempname)
                                for column in df.columns:
                                    for i in df.index:
                                        tempcontent.append(df[column][i])
                                searchtext3(str(tempcontent),text_filename)
                                df.Quit
                    shutil.rmtree(root+'/tempout')
                 text_filename.close()
                 searchtext(file_path+'_raw.txt',c,filetype)
            elif filename.endswith((".pdf")):
                 fullmatches = []
                 unique_matches = []
                 #page = 0
                 filetype = 'Adobe PDF'
                 req_image = []
                 full_text = []
                 file_path = os.path.join(root, filename)
                 if not os.path.exists('tmpimg'):
                     os.makedirs('tmpimg')
                 #with tempfile.TemporaryDirectory() as path:
                 images_from_path = convert_from_path(file_path, fmt='jpeg', output_folder='tmpimg')
                 for temproot, tempsubdirs, tempfiles in os.walk('tmpimg'):
                     for tempfile in tempfiles:
                         tempfile = 'tmpimg/' + tempfile
                         text = pytesseract.image_to_string(tempfile)
                         os.remove(tempfile)
                         searchtext2(text,c,filetype,filename,file_path,fullmatches,unique_matches)
                 if fullmatches:
                    match_report(filetype,filename,file_path,c,fullmatches)
                 else:
                    no_match_report(filetype,filename,file_path,c)
                 #pdfFileObj.close()
            else:
                 file_path = os.path.join(root, filename)
                 worksheet.write('A'+str(c),'Unknown')
                 worksheet.write('E'+str(c),filename)
                 worksheet.write('F'+str(c),file_path)
                 worksheet.write('G'+str(c),'Scanned')
                 worksheet.write('H'+str(c),'On File System')
                 worksheet.write('I'+str(c),'Out of Scope')
    os.remove(os.path.join(root, 'my-directory-list.txt'))
workbook.close()
